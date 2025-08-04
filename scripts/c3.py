# Generate: 1) PowerPoint deck for ApplicationAI, 2) Demo video script, 
# 3) vibefunder.ai landing page (HTML+CSS), and 4) a detailed planning doc + OpenAPI stub for agentic coding.
#
# If python-pptx is unavailable, fall back to a plain-text deck in Markdown.

import os, datetime, zipfile, json, textwrap, sys
from pathlib import Path

base = Path("./output/deliverables")
base.mkdir(parents=True, exist_ok=True)

today = datetime.datetime.now().strftime("%B %d, %Y")
law = "Massachusetts, USA"
start_date = datetime.date(2025, 8, 11)
m1_end = start_date + datetime.timedelta(weeks=4)
m2_end = m1_end + datetime.timedelta(weeks=6)
m3_end = m2_end + datetime.timedelta(weeks=5)

# -------- 1) Create PowerPoint (or Markdown fallback) --------
deck_path = base / "ApplicationAI_Charter_Deck.pptx"
fallback_md_path = base / "ApplicationAI_Charter_Deck.md"

slides = [
    {
        "title": "ApplicationAI — Charter Backer Briefing",
        "bullets": [f"Date: {today}", "Charter access • Code escrow • Milestone escrow (30/40/30)", "Audience: VC funds & innovation programs"],
        "notes": "Open by framing the pain: demo is easy, enterprise hard. We’re here to fund the hardening with 5 charter backers."
    },
    {
        "title": "Problem",
        "bullets": ["Inbound is messy: missing fields, scattered docs, inconsistent review", "Time-to-qualification is slow; reviewers lack standard briefs", "Security & compliance block adoption"],
        "notes": "Anchor on time wasted and risk. Mention that existing CRMs capture, they don’t evaluate."
    },
    {
        "title": "Solution",
        "bullets": ["URL → Application: auto-fill forms with citations", "Submitter copilot to resolve gaps, then approve", "Enrichment: market/comps/risks → explainable go/no-go brief"],
        "notes": "Transition to demo flow next."
    },
    {
        "title": "Demo Flow (10 min)",
        "bullets": ["1) URL→Application", "2) Submitter Copilot", "3) Enrichment with citations", "4) Rubric & export (Affinity/HubSpot)", "5) Security & Admin"],
        "notes": "Keep the demo tight and deterministic with sanitized data."
    },
    {
        "title": "What You Get as a Charter Backer",
        "bullets": ["Lifetime org license (or 5-Year All-Inclusive)", "Source-code escrow & continuity triggers", "Deployment options: SaaS / VPC / on‑prem", "Quarterly charter council; roadmap influence"],
        "notes": "Emphasize ‘buyer, not investor’: value and control without equity paperwork."
    },
    {
        "title": "Milestones & Dates",
        "bullets": [f"M1 Security & Identity — due {m1_end.isoformat()}", f"M2 Reliability & Data — due {m2_end.isoformat()}", f"M3 Compliance & Fit — due {m3_end.isoformat()}", "Funds released 30/40/30 on acceptance"],
        "notes": "Point to acceptance evidence and 10‑day review windows."
    },
    {
        "title": "Security & Compliance",
        "bullets": ["SSO (SAML/OIDC), RBAC, audit logging", "SBOM/SCA/DAST; pen test before M1 release", "DPA, SIG Lite/CAIQ responses, data residency", "No training on your data without explicit opt‑in"],
        "notes": "Offer to provide the short security packet immediately."
    },
    {
        "title": "Budget & Mechanics",
        "bullets": ["$100k total (5 × $20k) held in escrow", "Option: $5k/yr maintenance and security updates", "No equity, no rev‑share; commercial license + code escrow"],
        "notes": "Reiterate the simplicity of the commercial path vs. venture investment paperwork."
    },
    {
        "title": "Success Metrics",
        "bullets": ["<10 minutes to first qualified application", ">40% reviewer time saved / 100 apps", "AI pre‑score vs committee ≥ 0.75 (Spearman)", "Zero data leakage; ≤1 P0 incident/mo"],
        "notes": "These metrics become your proof points for broader GTM."
    },
    {
        "title": "Next Steps",
        "bullets": ["See the 15‑minute demo now", "Pick deployment (SaaS/VPC/on‑prem)", "Sign Charter License; funds deposit to escrow", "Week 1: charter council to lock v1 scope"],
        "notes": "Close decisively and book the follow‑up."
    }
]

def create_pptx(slides, path: Path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        bullet_layout = prs.slide_layouts[1]

        for i, s in enumerate(slides):
            layout = title_layout if i == 0 else bullet_layout
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = s["title"]
            if layout == bullet_layout:
                body = slide.placeholders[1].text_frame
                body.clear()
                for j, b in enumerate(s["bullets"]):
                    if j == 0:
                        body.text = b
                    else:
                        p = body.add_paragraph()
                        p.text = b
                        p.level = 0
            # Add notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = s["notes"]
        prs.save(path)
        return True
    except Exception as e:
        with open(fallback_md_path, "w") as f:
            for s in slides:
                f.write(f"# {s['title']}\n\n")
                for b in s["bullets"]:
                    f.write(f"- {b}\n")
                f.write(f"\n> Notes: {s['notes']}\n\n---\n\n")
        return False

ppt_ok = create_pptx(slides, deck_path)

# -------- 2) Demo video script --------
demo_script = f"""# ApplicationAI — Demo Video Script (8–10 minutes)

**Audience:** VC investment ops / platform teams  
**Goal:** Show URL→Application→Enrichment→Summary with security assurances and close on charter steps.

## 0) Cold open (0:00–0:20)
“Evaluating startups is slow because data is messy. ApplicationAI turns a company URL into a complete, explainable application in minutes.”

## 1) Setup (0:20–0:45)
- Show the landing screen with a URL field and “Target Program” selector.
- Mention: SSO is enabled; demo uses sanitized data.

## 2) URL → Application (0:45–2:30)
- Paste URL. Narrate: retrieval → extraction → field population with citations.
- Open the citations panel; click 2–3 sources to prove traceability.
- Point out missing required fields highlighted for submitter.

## 3) Submitter Copilot (2:30–4:00)
- Switch to the secure submitter link. The AI asks clarifying questions.
- Complete 2–3 fields. Show the “approve & lock” step.
- Narrate: all changes are logged; nothing is written to CRM until approval.

## 4) Enrichment (4:00–6:00)
- Open Market/Comps/Risks cards. Show explainability panel.
- Mention leakage/bias checks and deterministic re-run mode.
- Click “Generate Go/No-Go Brief.”

## 5) Rubric & Export (6:00–7:30)
- Show rubric auto-score; adjust one criterion to prove human-in-the-loop.
- Export to Affinity (demo org). Confirm via success toast and audit log entry.
- Mention CSV/Parquet export and admin console.

## 6) Security & Admin (7:30–8:30)
- Brief tour: SSO login screen, role settings, audit log, data export, and DPA link.
- “We never train foundation models on your data without opt-in.”

## 7) Close (8:30–10:00)
- Recap metrics: <10 min to a qualified application; >40% reviewer time saved.
- Charter mechanics: 5 backers × $20k = $100k in milestone escrow (30/40/30).
- “If you’d like to co-design v1, we can onboard your firm next week.”
"""

with open(base / "ApplicationAI_Demo_Script.md", "w") as f:
    f.write(demo_script)

# -------- 3) vibefunder.ai Landing Page (HTML + CSS) --------
html = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>VibeFunder — Ship the vibe. Not the pitch deck.</title>
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&display=swap" rel="stylesheet">
  <link rel="stylesheet" href="styles.css">
</head>
<body>
  <header class="nav">
    <div class="container">
      <div class="brand">VibeFunder<span class="dot">.ai</span></div>
      <nav>
        <a href="#how">How it works</a>
        <a href="#makers">Makers</a>
        <a href="#backers">Backers</a>
        <a href="#trust">Trust & Safety</a>
        <a href="#faq">FAQ</a>
      </nav>
    </div>
  </header>

  <section class="hero">
    <div class="container">
      <h1>Ship the vibe. Not the pitch deck.</h1>
      <p class="lead">VibeFunder is the alternative to venture capital for AI-native micro‑SaaS. Turn demos into dependable software with charter customers who fund the last mile—security, reliability, and compliance.</p>
      <form class="cta">
        <input type="email" placeholder="you@company.com" aria-label="Email">
        <button type="button">Request early access</button>
      </form>
      <div class="subtext">5% platform fee. Milestone escrow. Makers keep IP. Backers get lifetime licenses.</div>
    </div>
  </section>

  <section id="how" class="how">
    <div class="container grid">
      <div>
        <h2>How it works</h2>
        <ol>
          <li><strong>List your vibe-coded product.</strong> Demo video, scope, budget, and milestones.</li>
          <li><strong>Backers pledge.</strong> Funds held in escrow. No equity. Commercial license + code continuity.</li>
          <li><strong>Build with guardrails.</strong> Security, SSO, audit logs, SLAs via partner packs.</li>
          <li><strong>Release on acceptance.</strong> 30/40/30 escrow tied to evidence, not hype.</li>
        </ol>
      </div>
      <div class="card">
        <h3>Included rails</h3>
        <ul>
          <li>SSO/RBAC/audit logging baseline</li>
          <li>SOC2-track checklist + evidence vault</li>
          <li>Milestone acceptance workflows</li>
          <li>Source code escrow/mirror</li>
          <li>Partner marketplace (security, QA, SRE, legal)</li>
        </ul>
      </div>
    </div>
  </section>

  <section id="makers" class="split">
    <div class="container grid">
      <div>
        <h2>For Makers</h2>
        <p>Get upfront capital and a clear hardening checklist. Keep your IP. Choose SaaS, VPC, or on‑prem. Plug in partner packs for security and compliance.</p>
        <ul>
          <li>Charter lifetime or 5‑year licenses</li>
          <li>Transparent milestones & acceptance tests</li>
          <li>Vetted partners; you keep margin</li>
        </ul>
      </div>
      <div class="panel">Demo → Escrow → Enterprise Ready</div>
    </div>
  </section>

  <section id="backers" class="split alt">
    <div class="container grid">
      <div class="panel">Own your risk profile</div>
      <div>
        <h2>For Backers</h2>
        <p>Lock in lifetime value, shape the roadmap, and de‑risk with code continuity and compliance artifacts.</p>
        <ul>
          <li>Lifetime org license or on‑prem options</li>
          <li>Security & compliance packet (SIG Lite/CAIQ, DPA)</li>
          <li>Escrow tied to proof, not promises</li>
        </ul>
      </div>
    </div>
  </section>

  <section id="trust" class="trust">
    <div class="container">
      <h2>Trust & Safety</h2>
      <div class="grid">
        <div class="card">
          <h3>Verification</h3>
          <p>Makers screened; identity verified; repositories scanned.</p>
        </div>
        <div class="card">
          <h3>Continuity</h3>
          <p>Escrow/mirror triggers ensure code access for internal continuity.</p>
        </div>
        <div class="card">
          <h3>Evidence</h3>
          <p>Milestone acceptance requires artifacts—pen test summaries, SLO dashboards, restore drills.</p>
        </div>
      </div>
    </div>
  </section>

  <section id="faq" class="faq">
    <div class="container">
      <h2>FAQ</h2>
      <details>
        <summary>Is this crowdfunding or investing?</summary>
        <p>Neither. It’s a commercial pre‑purchase. Backers receive licenses and services, not equity or revenue share.</p>
      </details>
      <details>
        <summary>How does VibeFunder make money?</summary>
        <p>5% platform fee + margin on optional partner services (security, QA, SRE, legal).</p>
      </details>
      <details>
        <summary>Who owns the IP?</summary>
        <p>Makers own it. Backers get licenses and code‑continuity triggers.</p>
      </details>
      <details>
        <summary>What happens if a milestone fails?</summary>
        <p>There’s a cure window; otherwise escrow refunds release proportionally per terms.</p>
      </details>
    </div>
  </section>

  <footer class="footer">
    <div class="container">
      <span>© {datetime.datetime.now().year} VibeFunder</span>
      <a href="#">Terms</a>
      <a href="#">Privacy</a>
    </div>
  </footer>
</body>
</html>
"""

css = """*{box-sizing:border-box}body{margin:0;font-family:Inter,system-ui,-apple-system,Segoe UI,Roboto,Helvetica,Arial,sans-serif;color:#0b1320;background:#fcfcfe}
.container{max-width:1100px;margin:0 auto;padding:0 20px}
.nav{position:sticky;top:0;background:#ffffffcc;border-bottom:1px solid #eef;backdrop-filter:saturate(1.5) blur(6px)}
.nav .brand{font-weight:800;font-size:20px}.dot{color:#6757f5}
.nav nav a{margin-left:16px;color:#334;text-decoration:none}
.hero{padding:80px 0;background:linear-gradient(180deg,#f4f3ff,transparent)}
.hero h1{font-size:42px;margin:0 0 10px}
.lead{font-size:18px;max-width:820px}
.cta{margin-top:20px;display:flex;gap:10px}
.cta input{flex:1;padding:12px 14px;border:1px solid #ccd;border-radius:10px}
.cta button{padding:12px 18px;border:0;border-radius:10px;background:#6757f5;color:#fff;font-weight:600;cursor:pointer}
.subtext{margin-top:10px;color:#556}
.how,.split,.trust,.faq{padding:60px 0}
.grid{display:grid;grid-template-columns:1fr 1fr;gap:24px}
.card{background:#fff;border:1px solid #eef;border-radius:12px;padding:20px;box-shadow:0 2px 6px rgba(80,80,120,.06)}
.split .panel{display:flex;align-items:center;justify-content:center;border:1px dashed #cbd;height:180px;border-radius:12px;color:#556}
.split.alt{background:#f9f9ff}
.trust .grid{grid-template-columns:repeat(3,1fr)}
.faq details{margin:10px 0}.faq summary{cursor:pointer;font-weight:600}
.footer{padding:24px 0;border-top:1px solid #eef;color:#667;display:flex}
.footer .container{display:flex;justify-content:space-between;align-items:center}
@media(max-width:800px){.grid{grid-template-columns:1fr}.hero h1{font-size:34px}}
"""

lp_dir = base / "vibefunder_landing"
lp_dir.mkdir(exist_ok=True)
(lp_dir / "index.html").write_text(html)
(lp_dir / "styles.css").write_text(css)

# -------- 4) VibeFunder planning doc (for agentic coding) + OpenAPI stub --------
spec = f"""# VibeFunder.ai — Platform Planning Doc (Build‑Ready)

_Last updated {today}_

## 1) Vision
A marketplace where micro‑SaaS makers raise milestone‑escrowed funds from charter customers to harden prototypes into enterprise‑ready products. Platform takes 5% fee and margins on partner services.

## 2) Roles & Permissions
- **Maker (Founder):** create/edit campaign, connect repo, upload artifacts, request acceptance, book partner services.
- **Backer (Customer):** pledge funds, review evidence, accept/holdback/reject milestones, access license & escrow terms.
- **Partner (Vendor):** list services, accept work orders, upload deliverables, invoice via platform.
- **Platform Admin:** KYC/verification, dispute board ops, fee config, partner curation.
- **Auditor (read‑only):** view artifacts/notes for compliance (optional).

## 3) Core Objects
- **Campaign**: title, problem, solution, video, budget, milestones[], badges[], status, repo_url, deploy_modes[]
- **Milestone**: name, % payout, acceptance_criteria, evidence[], status, due_date
- **Pledge**: backer_id, campaign_id, amount, status (authorized|captured|refunded)
- **Escrow**: wallet_id, balances (m1/m2/m3), release_rules
- **Artifact**: type (pen_test, sbom, demo_video, slo_screenshot, restore_log), url, checksum
- **Badge**: code (SECURITY_READY, SOC2_TRACK, SSO_READY, PII_SAFE), status
- **PartnerService**: category, rate_card, sla, sample_deliverables
- **WorkOrder**: partner_id, campaign_id, scope, price, status, deliverables[]
- **License**: type (Lifetime, 5Year, OnPrem), terms_url, escrow_triggers
- **User**: role(s), org, idp (SAML/OIDC), verification_status

## 4) Key Flows (Happy Paths)
1. Maker creates Campaign → uploads demo → defines milestones & acceptance tests → requests listing review.  
2. Platform verifies Maker (KYC) → approves listing → campaign goes live.  
3. Backer pledges → funds authorized → when fully funded, M0 starts.  
4. Maker completes M1 → uploads evidence → Backers review → accept or holdback → escrow releases 30%.  
5. Repeat for M2 (40%) and M3 (30%).  
6. On completion, licenses issued; partners invoice; badges displayed.

## 5) Non‑Happy Paths
- Milestone rejected → cure window → re‑review → proportional refund if unresolved.  
- Maker default (abandon) → dispute board → partial refund + code escrow trigger.  
- Backer chargeback → freeze releases; initiate dispute.

## 6) Data Model (DDL sketch)
```sql
-- Postgres
create table users(id uuid primary key, email text unique, org_id uuid, roles text[], idp_sub text, created_at timestamptz default now());
create table campaigns(id uuid primary key, maker_id uuid references users(id), title text, summary text, budget_cents int,
  currency text default 'USD', status text, repo_url text, video_url text, deploy_modes text[], created_at timestamptz default now());
create table milestones(id uuid primary key, campaign_id uuid references campaigns(id), name text, pct int, due_date date,
  acceptance_json jsonb, status text, evidence jsonb default '[]'::jsonb);
create table pledges(id uuid primary key, campaign_id uuid references campaigns(id), backer_id uuid references users(id),
  amount_cents int, status text, payment_ref text, created_at timestamptz default now());
create table escrows(id uuid primary key, campaign_id uuid references campaigns(id), balances jsonb, rules jsonb);
create table artifacts(id uuid primary key, campaign_id uuid references campaigns(id), kind text, url text, checksum text, meta jsonb);
create table badges(id uuid primary key, campaign_id uuid references campaigns(id), code text, status text, evidence jsonb);
create table partner_services(id uuid primary key, partner_id uuid references users(id), category text, rate_card jsonb, sla jsonb);
create table work_orders(id uuid primary key, campaign_id uuid, partner_id uuid, scope jsonb, price_cents int, status text, deliverables jsonb);
create table licenses(id uuid primary key, campaign_id uuid references campaigns(id), type text, terms_url text, triggers jsonb);
```

## 7) API (REST, JWT + SSO; OpenAPI stub provided)
- Auth: OIDC login → JWT. Backers can also use API keys for finance operations (limited).  
- Idempotency‑keys on pledge & release endpoints.  
- Webhooks for milestone submissions, acceptance, and refund events.

## 8) Escrow & Payments
- Provider: Stripe (PaymentIntents, Connect, custom accounts) or Adyen.  
- Funds authorized at pledge; captured at campaign close; held in platform-owned segregated account with ledger for M1/M2/M3.  
- Release operations create transfers less 5% platform fee; partial refunds pro‑rata.

## 9) Compliance
- KYC/KYB on Makers & Partners.  
- Data Processing Addendum templates; subprocessor registry.  
- Optional SOC2‑Track evidence vault (files + attestations).

## 10) Architecture (suggested)
- **Frontend:** Next.js (App Router), TypeScript, React Server Components.  
- **Backend:** NestJS (Node), Postgres, Prisma, Redis (queues), OpenAPI.  
- **Storage:** S3-compatible for artifacts; signed URLs.  
- **Auth:** Auth0/WorkOS/Keycloak for SAML/OIDC.  
- **Infra:** AWS (ECS Fargate or Lambda), CloudFront, WAF, KMS.  
- **Observability:** OpenTelemetry, Prometheus, Loki; Alertmanager.  
- **AI modules:** Python microservices (FastAPI) for document parsing & scoring (optional).

## 11) Events (for agents)
- campaign.created, campaign.approved, pledge.authorized, pledge.captured, milestone.submitted, milestone.accepted, milestone.rejected, escrow.released, escrow.refunded, workorder.created, workorder.completed.

## 12) Acceptance Tests (platform)
- Create campaign → approve → fund → submit M1 → accept → auto-release 30% → ledger matches.  
- Reject with holdback → cure → accept → release remainder.  
- Dispute path exercises refund logic & escrow ledger.

## 13) Backlog (MVP → v1)
- MVP: Auth (OIDC), Campaign CRUD, Pledge/Checkout, Escrow Ledger, Milestone Workflow, Artifact Upload, Badges, Admin panel.  
- v1: Partner Marketplace, Work Orders & Invoicing, Dispute Board, Public badges, Analytics, Email notifications.

## 14) Security
- Role-based permissions; signed URLs; row-level security where possible.  
- Rate limits & WAF; audit logging for sensitive ops.  
- SBOM, SCA/DAST in CI; secrets scanning; dependency pinning.

## 15) GTM
- Seed with 3–5 flagship campaigns.  
- Publish milestone evidence & uptime stats; build trust.  
- Partner with 2–3 security vendors; revenue share.

"""

openapi = {
  "openapi": "3.0.3",
  "info": {"title": "VibeFunder API", "version": "0.1.0"},
  "servers": [{"url": "https://api.vibefunder.ai"}],
  "paths": {
    "/auth/login": {"post": {"summary": "OIDC callback", "responses": {"200": {"description": "OK"}}}},
    "/campaigns": {
      "get": {"summary": "List campaigns", "responses": {"200": {"description": "OK"}}},
      "post": {"summary": "Create campaign", "responses": {"201": {"description": "Created"}}}
    },
    "/campaigns/{id}": {
      "get": {"summary": "Get campaign", "responses": {"200": {"description": "OK"}}},
      "patch": {"summary": "Update campaign", "responses": {"200": {"description": "OK"}}}
    },
    "/campaigns/{id}/milestones": {
      "get": {"summary": "List milestones", "responses": {"200": {"description": "OK"}}},
      "post": {"summary": "Create milestone", "responses": {"201": {"description": "Created"}}}
    },
    "/milestones/{id}/submit": {
      "post": {"summary": "Submit milestone for review", "responses": {"200": {"description": "Submitted"}}}
    },
    "/milestones/{id}/accept": {
      "post": {"summary": "Accept milestone", "responses": {"200": {"description": "Accepted"}}}
    },
    "/milestones/{id}/reject": {
      "post": {"summary": "Reject milestone", "responses": {"200": {"description": "Rejected"}}}
    },
    "/pledges": {
      "post": {"summary": "Create pledge (authorize payment)", "responses": {"201": {"description": "Created"}}}
    },
    "/escrow/{campaignId}/release": {
      "post": {"summary": "Release escrow for a milestone", "responses": {"200": {"description": "Released"}}}
    },
    "/artifacts": {
      "post": {"summary": "Upload artifact metadata", "responses": {"201": {"description": "Created"}}}
    },
    "/badges/{campaignId}": {
      "get": {"summary": "List badges", "responses": {"200": {"description": "OK"}}}
    }
  },
  "components": {}
}

(base / "vibefunder_platform_spec.md").write_text(spec)
(base / "vibefunder_openapi.json").write_text(json.dumps(openapi, indent=2))

# Bundle everything in a single zip
zip_path = "./output/VibeFunder_and_ApplicationAI_Pack.zip"
with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:
    for p in [deck_path if deck_path.exists() else fallback_md_path,
              base / "ApplicationAI_Demo_Script.md",
              lp_dir / "index.html",
              lp_dir / "styles.css",
              base / "vibefunder_platform_spec.md",
              base / "vibefunder_openapi.json"]:
        z.write(p, arcname=p.name)

zip_path, deck_path.exists(), str(deck_path)